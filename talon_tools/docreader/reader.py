"""Document reader — extracts text from PDF, Excel, Word, PowerPoint."""

from __future__ import annotations

import io
import logging
from pathlib import Path

log = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".xlsx", ".xls", ".docx", ".pptx"}


def read_document(path: str | Path) -> str:
    """Read a document and return its text content.

    Dispatches to the appropriate reader based on file extension.
    """
    p = Path(path)
    ext = p.suffix.lower()

    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    if ext == ".pdf":
        return _read_pdf(p)
    elif ext in (".xlsx", ".xls"):
        return _read_excel(p)
    elif ext == ".docx":
        return _read_docx(p)
    elif ext == ".pptx":
        return _read_pptx(p)
    else:
        raise ValueError(f"Unsupported format: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")


def read_document_bytes(data: bytes, filename: str) -> str:
    """Read a document from bytes (e.g. downloaded file)."""
    ext = Path(filename).suffix.lower()
    buf = io.BytesIO(data)

    if ext == ".pdf":
        return _read_pdf_stream(buf)
    elif ext in (".xlsx", ".xls"):
        return _read_excel_stream(buf)
    elif ext == ".docx":
        return _read_docx_stream(buf)
    elif ext == ".pptx":
        return _read_pptx_stream(buf)
    else:
        raise ValueError(f"Unsupported format: {ext}")


# -- PDF ---------------------------------------------------------------------

def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return _extract_pdf(reader)


def _read_pdf_stream(buf: io.BytesIO) -> str:
    from pypdf import PdfReader
    reader = PdfReader(buf)
    return _extract_pdf(reader)


def _extract_pdf(reader) -> str:
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
    return "\n\n".join(pages) if pages else "(No text content found in PDF)"


# -- Excel -------------------------------------------------------------------

def _read_excel(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    return _extract_excel(wb)


def _read_excel_stream(buf: io.BytesIO) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(buf, read_only=True, data_only=True)
    return _extract_excel(wb)


def _extract_excel(wb) -> str:
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            header = rows[0]
            sep = " | ".join("---" for _ in rows[0].split(" | "))
            table = "\n".join([header, sep] + rows[1:])
            sheets.append(f"## Sheet: {name}\n\n{table}")
    wb.close()
    return "\n\n".join(sheets) if sheets else "(No data found in workbook)"


# -- Word (docx) -------------------------------------------------------------

def _read_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return _extract_docx(doc)


def _read_docx_stream(buf: io.BytesIO) -> str:
    from docx import Document
    doc = Document(buf)
    return _extract_docx(doc)


def _extract_docx(doc) -> str:
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if style.startswith("Heading 1"):
            parts.append(f"# {text}")
        elif style.startswith("Heading 2"):
            parts.append(f"## {text}")
        elif style.startswith("Heading 3"):
            parts.append(f"### {text}")
        elif style.startswith("List"):
            parts.append(f"- {text}")
        else:
            parts.append(text)

    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            header = rows[0]
            sep = " | ".join("---" for _ in rows[0].split(" | "))
            parts.append("\n".join([header, sep] + rows[1:]))

    return "\n\n".join(parts) if parts else "(No text content found in document)"


# -- PowerPoint --------------------------------------------------------------

def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    return _extract_pptx(prs)


def _read_pptx_stream(buf: io.BytesIO) -> str:
    from pptx import Presentation
    prs = Presentation(buf)
    return _extract_pptx(prs)


def _extract_pptx(prs) -> str:
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "(No text content found in presentation)"
